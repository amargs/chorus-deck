from pptx import Presentation

TITLE_PREFIX = "Cântico:"

def read_ppt(file_path: str):
    ppt = Presentation(file_path)
    slide_titles = []

    for slide in ppt.slides:
        title_shape = slide.shapes.title
        if title_shape and title_shape.text:
            slide_titles.append(clean_title(title_shape.text)) 

    return slide_titles

def clean_title(title: str) -> str:
    return title.replace(TITLE_PREFIX, "").strip()
