from copy import deepcopy
from pptx import Presentation
from .models import Song

TITLE_PREFIX = "Cântico:"

def read_song_titles(source_file_path: str) -> list[str]:
    ppt = Presentation(source_file_path)
    slide_titles = []

    for slide in ppt.slides:
        title_shape = slide.shapes.title
        if title_shape and title_shape.text:
            slide_titles.append(clean_title(title_shape.text))

    return slide_titles

def index_songs(source_file_path: str) -> list[Song]:
    ppt = Presentation(source_file_path)
    song_id = 1
    song_data = []

    for i, slide in enumerate(ppt.slides, start=0):
        title_shape = slide.shapes.title

        if title_shape and title_shape.text:
            song_data.append(Song(id=song_id, title=clean_title(title_shape.text), slide_start=i, slide_end=i))
            song_id += 1
        elif slide_is_empty(slide) and len(song_data) > 0:
            song_data[-1].slide_end = i-1
        
    if len(song_data) > 0:
        song_data[-1].slide_end = len(ppt.slides) - 1

    return 

def create_ppt(source_file_path: str, songs: list[Song]) -> Presentation:
    slide_ranges = []
    offset = 0

    for song in sorted(songs, key=lambda x: x.id):
        slide_ranges.append((song.slide_start, song.slide_end))
        song.slide_end = offset + (song.slide_end - song.slide_start)
        song.slide_start = offset
        offset = song.slide_end + 1

    ppt = delete_unwanted_slides(source_file_path, slide_ranges)

    ordered_slides = []
    for song in songs:
        ordered_slides.extend(range(song.slide_start, song.slide_end + 1))

    reorder_slides(ppt, ordered_slides)

    return ppt

def delete_unwanted_slides(source_file_path: str, slide_ranges: list[tuple[int, int]]) -> Presentation:
    ppt = Presentation(source_file_path)

    slide_idx_to_keep = {i for start, end in slide_ranges for i in range(start, end + 1)}

    for i in reversed(range(len(ppt.slides))):
        if i not in slide_idx_to_keep:
            slide_id = ppt.slides._sldIdLst[i].rId
            ppt.part.drop_rel(slide_id)
            del ppt.slides._sldIdLst[i]

    return ppt

def reorder_slides(ppt: Presentation, ordered_slides: list[int]):
    slide_id_list = ppt.slides._sldIdLst
    reordered_slide_id_list = [slide_id_list[idx] for idx in ordered_slides]

    for id in slide_id_list:
        slide_id_list.remove(id)

    for id in reordered_slide_id_list:
        slide_id_list.append(id)

def clean_title(title: str) -> str:
    return title.replace(TITLE_PREFIX, "").strip()

def slide_is_empty(slide) -> bool:
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text.strip():
            return False
    return True
